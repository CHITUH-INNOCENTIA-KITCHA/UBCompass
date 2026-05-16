#!/usr/bin/env python3
"""
UBCompass PowerPoint Presentation Generator
Creates a professional, well-designed presentation with animations
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap, qn
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor
import os

# Alias for convenience
RgbColor = RGBColor

# UB Brand Colors
UB_GREEN = RgbColor(0, 100, 0)  # #006400
UB_GREEN_LIGHT = RgbColor(168, 213, 162)  # #A8D5A2
UB_GREEN_SOFT = RgbColor(246, 251, 244)  # #F6FBF4
WHITE = RgbColor(255, 255, 255)
DARK_TEXT = RgbColor(26, 26, 26)
MUTED_TEXT = RgbColor(100, 100, 100)
ORANGE = RgbColor(255, 107, 53)  # Route color
BLUE = RgbColor(66, 133, 244)  # GPS blue

def add_entrance_animation(shape, animation_type="fade"):
    """Add entrance animation to a shape (note: limited support in python-pptx)"""
    # python-pptx has limited animation support, but we can add transition effects to slides
    pass

def create_title_slide(prs, title, subtitle):
    """Create a title slide with UB branding"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add green header band
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3))
    header.fill.solid()
    header.fill.fore_color.rgb = UB_GREEN
    header.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = UB_GREEN_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Add decorative element
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(3.5), Inches(1.6), Inches(1.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = UB_GREEN
    circle.line.fill.background()

    # Add compass icon text
    icon_box = slide.shapes.add_textbox(Inches(4.2), Inches(3.85), Inches(1.6), Inches(1))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🧭"
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_section_slide(prs, section_number, section_title, icon="📍"):
    """Create a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full green background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = UB_GREEN
    bg.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{icon}  Section {section_number}"
    p.font.size = Pt(24)
    p.font.color.rgb = UB_GREEN_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, title, bullets, has_icon=True):
    """Create a standard content slide with bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UB_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    return slide

def create_two_column_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    """Create a two-column layout slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UB_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column title
    if left_title:
        left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5))
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = UB_GREEN

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)

    # Right column title
    if right_title:
        right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5))
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = UB_GREEN

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)

    return slide

def create_table_slide(prs, title, headers, rows):
    """Create a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UB_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * num_rows)).table

    # Style header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = UB_GREEN
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER

    # Fill data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = UB_GREEN_SOFT

    return slide

def create_architecture_slide(prs):
    """Create the system architecture slide with diagram"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UB_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Main app box
    app_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(1.5), Inches(6), Inches(1))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = UB_GREEN
    app_box.line.color.rgb = UB_GREEN

    app_text = slide.shapes.add_textbox(Inches(2), Inches(1.7), Inches(6), Inches(0.6))
    tf = app_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📱 UBCompass Mobile App (React Native + Expo)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Four tab boxes
    tabs = [("🗺️ Map", 0.3), ("🔍 Search", 2.55), ("🏢 Indoor", 4.8), ("⚙️ Settings", 7.05)]
    for tab_name, x_pos in tabs:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(2.7), Inches(2.1), Inches(0.6))
        box.fill.solid()
        box.fill.fore_color.rgb = UB_GREEN_LIGHT
        box.line.color.rgb = UB_GREEN

        text = slide.shapes.add_textbox(Inches(x_pos), Inches(2.8), Inches(2.1), Inches(0.5))
        tf = text.text_frame
        p = tf.paragraphs[0]
        p.text = tab_name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER

    # State management box
    state_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3.5), Inches(7), Inches(0.5))
    state_box.fill.solid()
    state_box.fill.fore_color.rgb = RgbColor(230, 230, 230)
    state_box.line.color.rgb = UB_GREEN

    state_text = slide.shapes.add_textbox(Inches(1.5), Inches(3.55), Inches(7), Inches(0.5))
    tf = state_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📦 State Management (Zustand)"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # External services
    services = [
        ("📍 Expo\nLocation", 0.3, RgbColor(66, 133, 244)),
        ("🗄️ Supabase\nDatabase", 2.55, RgbColor(62, 207, 142)),
        ("🛣️ OSRM\nRouting", 4.8, ORANGE),
        ("🗺️ OpenStreet\nMap", 7.05, RgbColor(130, 180, 100))
    ]

    for service_name, x_pos, color in services:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(4.5), Inches(2.1), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        text = slide.shapes.add_textbox(Inches(x_pos), Inches(4.65), Inches(2.1), Inches(0.8))
        tf = text.text_frame
        p = tf.paragraphs[0]
        p.text = service_name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Arrows (simple lines)
    for x in [1.35, 3.6, 5.85, 8.1]:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.1), Inches(0.05), Inches(0.35))
        line.fill.solid()
        line.fill.fore_color.rgb = MUTED_TEXT
        line.line.fill.background()

    return slide

def create_feature_highlight_slide(prs, title, feature_icon, description, points):
    """Create a feature highlight slide with icon"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UB_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Feature icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(1.5), Inches(1.5), Inches(1.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = UB_GREEN_LIGHT
    circle.line.color.rgb = UB_GREEN
    circle.line.width = Pt(3)

    icon_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(1.5), Inches(1))
    tf = icon_text.text_frame
    p = tf.paragraphs[0]
    p.text = feature_icon
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(Inches(2.3), Inches(1.5), Inches(7), Inches(1))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_TEXT

    # Feature points
    points_box = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(3.5))
    tf = points_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(10)

    return slide

def create_demo_slide(prs, step_num, step_title, description, duration):
    """Create a demo step slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = UB_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Demo Step {step_num}: {step_title}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4), Inches(1.8), Inches(2), Inches(2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = UB_GREEN
    circle.line.fill.background()

    num_text = slide.shapes.add_textbox(Inches(4), Inches(2.3), Inches(2), Inches(1.2))
    tf = num_text.text_frame
    p = tf.paragraphs[0]
    p.text = str(step_num)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Duration badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5.8), Inches(2), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.fill.background()

    dur_text = slide.shapes.add_textbox(Inches(4), Inches(5.88), Inches(2), Inches(0.4))
    tf = dur_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"⏱️ {duration}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_conclusion_slide(prs):
    """Create the conclusion slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full green background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = UB_GREEN
    bg.line.fill.background()

    # Compass icon
    icon_text = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(2), Inches(1.5))
    tf = icon_text.text_frame
    p = tf.paragraphs[0]
    p.text = "🧭"
    p.font.size = Pt(80)
    p.alignment = PP_ALIGN.CENTER

    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "UBCompass - Navigating the Future of Campus Mobility"
    p.font.size = Pt(24)
    p.font.color.rgb = UB_GREEN_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Credits
    credits = [
        "Supervisor: Dr. Mougnol Romeo",
        "Development Team: Group 5",
        "University of Buea, Cameroon"
    ]

    for i, credit in enumerate(credits):
        credit_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5 + i * 0.5), Inches(9), Inches(0.5))
        tf = credit_box.text_frame
        p = tf.paragraphs[0]
        p.text = credit
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: Title Slide =====
    create_title_slide(
        prs,
        "UBCompass",
        "Campus Navigation Mobile Application for University of Buea"
    )

    # ===== SLIDE 2: Agenda =====
    create_content_slide(prs, "📋 Presentation Agenda", [
        "Introduction & Problem Statement",
        "Solution Overview - What is UBCompass?",
        "Target Audience & Use Cases",
        "System Architecture",
        "Technology Stack",
        "Key Features Deep Dive",
        "Database Design",
        "Live Demo",
        "Future Roadmap",
        "Q&A Session"
    ])

    # ===== SLIDE 3: Section - Introduction =====
    create_section_slide(prs, "1", "Introduction", "🎯")

    # ===== SLIDE 4: The Problem =====
    create_content_slide(prs, "🚨 The Problem", [
        "New students struggle to find faculties, lecture halls, and offices",
        "Visitors have difficulty navigating the large UB campus",
        "No existing digital navigation solution for the campus",
        "Indoor navigation is virtually non-existent",
        "Students with disabilities need wheelchair-accessible routes",
        "Paper maps are outdated and hard to use"
    ])

    # ===== SLIDE 5: The Solution =====
    create_content_slide(prs, "💡 The Solution: UBCompass", [
        "Cross-platform mobile app (iOS & Android)",
        "Real-time GPS-based outdoor navigation",
        "Turn-by-turn walking directions between any two points",
        "Indoor floor plans with room-level navigation",
        "Accessibility-aware routing (wheelchair-friendly paths)",
        "Building and room search functionality",
        "Accurate walking time and distance calculations"
    ])

    # ===== SLIDE 6: Section - Users =====
    create_section_slide(prs, "2", "Target Audience", "👥")

    # ===== SLIDE 7: User Types Table =====
    create_table_slide(prs, "👥 Who Uses UBCompass?",
        ["User Type", "Primary Need", "Key Feature"],
        [
            ["New Students", "Find classrooms & faculties", "GPS Directions"],
            ["Visitors", "Locate admin buildings", "Building Search"],
            ["Disabled Users", "Wheelchair routes", "Accessibility Mode"],
            ["Staff Members", "Find offices & rooms", "Indoor Navigation"],
            ["Event Attendees", "Locate venues", "Building Details"]
        ]
    )

    # ===== SLIDE 8: Use Case Example =====
    create_content_slide(prs, "📖 Use Case: First Day at UB", [
        "Scenario: A freshman arrives at UB Main Gate",
        "Goal: Find the Faculty of Science (FOS)",
        "Step 1: Open UBCompass app",
        "Step 2: Search for 'Faculty of Science'",
        "Step 3: Tap 'Get Directions'",
        "Step 4: Follow the blue route on the map",
        "Step 5: See live GPS position (blue dot)",
        "Result: Arrives at FOS in 5 minutes! ✓"
    ])

    # ===== SLIDE 9: Section - Architecture =====
    create_section_slide(prs, "3", "System Architecture", "🏗️")

    # ===== SLIDE 10: Architecture Diagram =====
    create_architecture_slide(prs)

    # ===== SLIDE 11: Data Flow =====
    create_content_slide(prs, "🔄 Data Flow", [
        "User Action → React Component → Custom Hook → External Service",
        "",
        "Example: Getting Directions",
        "1. User taps 'Get Directions' button",
        "2. DirectionsScreen component triggered",
        "3. useLocation hook provides GPS coordinates",
        "4. osmRouting utility calls OSRM API",
        "5. Route coordinates returned to state",
        "6. MapView renders route as Polyline",
        "7. ETA calculated using Haversine formula"
    ])

    # ===== SLIDE 12: Section - Tech Stack =====
    create_section_slide(prs, "4", "Technology Stack", "🛠️")

    # ===== SLIDE 13: Frontend Technologies =====
    create_table_slide(prs, "🎨 Frontend Technologies",
        ["Technology", "Version", "Purpose"],
        [
            ["React Native", "0.81.5", "Cross-platform framework"],
            ["Expo", "54.0", "Development platform"],
            ["TypeScript", "5.9.2", "Type-safe JavaScript"],
            ["React Native Paper", "5.15.1", "Material Design 3 UI"],
            ["React Native SVG", "15.12.1", "Floor plan rendering"]
        ]
    )

    # ===== SLIDE 14: Backend & APIs =====
    create_table_slide(prs, "🌐 Backend & External APIs",
        ["Service", "Type", "Purpose"],
        [
            ["Supabase", "BaaS", "PostgreSQL database"],
            ["OpenStreetMap", "Map Tiles", "Outdoor map rendering"],
            ["OSRM", "Routing API", "Walking directions"],
            ["Expo Location", "Device API", "GPS tracking"],
            ["Zustand", "State Mgmt", "Global app state"]
        ]
    )

    # ===== SLIDE 15: Section - Features =====
    create_section_slide(prs, "5", "Key Features", "⭐")

    # ===== SLIDE 16: Feature 1 - GPS Navigation =====
    create_feature_highlight_slide(prs,
        "Feature 1: GPS Navigation",
        "📍",
        "Real-time GPS tracking shows user position as a 'blue dot' on the map with walking directions to any campus building.",
        [
            "Live GPS tracking with accuracy circle",
            "OSRM API for actual walking paths (not straight lines)",
            "Turn-by-turn route visualization",
            "Distance and ETA calculations",
            "Works with any starting point"
        ]
    )

    # ===== SLIDE 17: Feature 2 - Indoor Navigation =====
    create_feature_highlight_slide(prs,
        "Feature 2: Indoor Navigation",
        "🏢",
        "SVG-based floor plans allow navigation inside buildings with room-level directions using Dijkstra's shortest path algorithm.",
        [
            "Interactive SVG floor plans",
            "Pinch-to-zoom and pan gestures",
            "Room tap selection",
            "Floor switching (G, 1, 2, etc.)",
            "Dijkstra's algorithm for shortest path"
        ]
    )

    # ===== SLIDE 18: Feature 3 - Accessibility =====
    create_feature_highlight_slide(prs,
        "Feature 3: Accessibility Mode",
        "♿",
        "When enabled, all routes automatically avoid stairs and prefer ramps and accessible entrances.",
        [
            "Global toggle in Settings",
            "Indoor routing avoids stairs nodes",
            "Shows accessible entrance info",
            "Ramp-only floor transitions",
            "Visual indicator when active"
        ]
    )

    # ===== SLIDE 19: Feature 4 - Search =====
    create_feature_highlight_slide(prs,
        "Feature 4: Smart Search",
        "🔍",
        "Instant search across all campus buildings and rooms with category filtering and recent search history.",
        [
            "Real-time filtering as you type",
            "Search by building name or room",
            "Category icons (faculty, admin, facility)",
            "Recent searches saved",
            "Direct navigation from results"
        ]
    )

    # ===== SLIDE 20: Feature 5 - ETA Calculation =====
    create_feature_highlight_slide(prs,
        "Feature 5: Accurate ETA",
        "⏱️",
        "Walking time estimates use the Haversine formula for distance and realistic walking speeds.",
        [
            "Haversine formula accounts for Earth's curvature",
            "Default walking speed: 5 km/h",
            "Example: 500m ≈ 6 minutes",
            "Example: 1 km ≈ 12 minutes",
            "Updates based on actual route distance"
        ]
    )

    # ===== SLIDE 21: Section - Database =====
    create_section_slide(prs, "6", "Database Design", "🗄️")

    # ===== SLIDE 22: Database Schema =====
    create_two_column_slide(prs, "📊 PostgreSQL Schema",
        [
            "buildings table:",
            "  - id (UUID, primary key)",
            "  - name, short_name",
            "  - latitude, longitude",
            "  - floors, has_indoor_map",
            "  - category, image_url",
            "  - accessible_entrance",
            "  - opening_hours"
        ],
        [
            "rooms table:",
            "  - id (UUID, primary key)",
            "  - building_id (foreign key)",
            "  - name, room_number",
            "  - floor, type",
            "  - node_x, node_y",
            "  - accessible"
        ],
        "Buildings Table",
        "Rooms Table"
    )

    # ===== SLIDE 23: Campus Data =====
    create_table_slide(prs, "🏛️ Seeded Campus Buildings",
        ["Building", "Category", "Coordinates"],
        [
            ["Main Gate", "Facility", "4.1537, 9.2837"],
            ["Faculty of Science", "Faculty", "4.1549, 9.2841"],
            ["Faculty of Engineering", "Faculty", "4.1558, 9.2853"],
            ["Administration Block", "Admin", "4.1540, 9.2855"],
            ["University Library", "Facility", "4.1545, 9.2860"],
            ["Sports Complex", "Facility", "4.1520, 9.2840"]
        ]
    )

    # ===== SLIDE 24: Section - Algorithms =====
    create_section_slide(prs, "7", "Algorithms", "🧮")

    # ===== SLIDE 25: Haversine Formula =====
    create_content_slide(prs, "📐 Haversine Formula", [
        "Purpose: Calculate distance between two GPS coordinates",
        "",
        "Formula:",
        "  a = sin²(Δφ/2) + cos(φ1) × cos(φ2) × sin²(Δλ/2)",
        "  c = 2 × atan2(√a, √(1−a))",
        "  d = R × c",
        "",
        "Where:",
        "  R = Earth's radius (6,371 km)",
        "  φ = latitude in radians",
        "  λ = longitude in radians"
    ])

    # ===== SLIDE 26: Dijkstra's Algorithm =====
    create_content_slide(prs, "🛤️ Dijkstra's Shortest Path", [
        "Purpose: Find shortest path between rooms indoors",
        "",
        "How it works:",
        "1. Build graph from room connections",
        "2. Initialize: start = 0, all others = ∞",
        "3. Process nodes by shortest distance",
        "4. Update neighbor distances if shorter path found",
        "5. Continue until destination reached",
        "6. Backtrack to reconstruct path",
        "",
        "Accessibility Mode: Filters out 'stairs' nodes"
    ])

    # ===== SLIDE 27: Section - Demo =====
    create_section_slide(prs, "8", "Live Demo", "🎬")

    # ===== SLIDE 28-33: Demo Steps =====
    demo_steps = [
        ("App Launch", "Show splash screen and initial map view", "30 sec"),
        ("GPS Permission", "Grant location, show blue dot appearing", "30 sec"),
        ("Search Building", "Search 'Faculty of Science', view details", "1 min"),
        ("Get Directions", "Show route, distance chip, ETA chip", "1 min"),
        ("Indoor Navigation", "Select floor, find room, show path", "2 min"),
        ("Accessibility Mode", "Toggle on, show route avoiding stairs", "30 sec")
    ]

    for i, (title, desc, duration) in enumerate(demo_steps, 1):
        create_demo_slide(prs, i, title, desc, duration)

    # ===== SLIDE 34: Section - Future =====
    create_section_slide(prs, "9", "Future Roadmap", "🚀")

    # ===== SLIDE 35: Future Features =====
    create_two_column_slide(prs, "🔮 Future Enhancements",
        [
            "Short-term:",
            "More indoor maps (FET, FSMS)",
            "Offline map caching",
            "Campus events integration",
            "Voice navigation"
        ],
        [
            "Long-term:",
            "Augmented Reality (AR) overlay",
            "Real-time bus tracking",
            "Crowd density indicators",
            "Multi-language support"
        ],
        "Coming Soon",
        "Future Vision"
    )

    # ===== SLIDE 36: Development Phases =====
    create_table_slide(prs, "📈 Development Progress",
        ["Phase", "Description", "Status"],
        [
            ["Phase 1", "Project Setup & Structure", "✅ 100%"],
            ["Phase 2", "Supabase Database", "✅ 100%"],
            ["Phase 3", "Outdoor Navigation", "✅ 100%"],
            ["Phase 4", "Search & Discovery", "✅ 100%"],
            ["Phase 5", "Indoor Navigation", "✅ 100%"],
            ["Phase 6", "Settings & Accessibility", "✅ 100%"],
            ["Phase 7", "Polish & Integration", "✅ 100%"],
            ["Phase 8", "Testing & Deployment", "🔄 In Progress"]
        ]
    )

    # ===== SLIDE 37: Key Achievements =====
    create_content_slide(prs, "🏆 Key Achievements", [
        "Cross-platform: Single codebase for iOS and Android",
        "Real-world data: Actual UB campus coordinates",
        "Accessibility-first: Wheelchair-friendly routing",
        "Open-source stack: No vendor lock-in, completely free",
        "Modern UI: Material Design 3 compliance",
        "Accurate navigation: GPS + OSRM routing",
        "Indoor innovation: SVG floor plans with pathfinding"
    ])

    # ===== SLIDE 38: Conclusion =====
    create_conclusion_slide(prs)

    # ===== SLIDE 39: Q&A =====
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full green background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = UB_GREEN
    bg.line.fill.background()

    # Q&A text
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Answers"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "❓ We're happy to answer your questions!"
    p.font.size = Pt(24)
    p.font.color.rgb = UB_GREEN_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Save the presentation
    output_path = "/home/daytona/project/UBCompass_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return output_path

if __name__ == "__main__":
    main()
