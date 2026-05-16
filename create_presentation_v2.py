#!/usr/bin/env python3
"""
UBCompass PowerPoint Presentation Generator v2
Creates a modern, professional presentation with better design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# Modern Color Palette
PRIMARY = RGBColor(0, 100, 0)       # UB Green #006400
PRIMARY_DARK = RGBColor(0, 70, 0)   # Darker green
PRIMARY_LIGHT = RGBColor(200, 230, 200)  # Light green
ACCENT = RGBColor(255, 107, 53)     # Orange accent
ACCENT_BLUE = RGBColor(66, 133, 244)  # Blue for GPS
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(33, 33, 33)
GRAY = RGBColor(120, 120, 120)
LIGHT_BG = RGBColor(248, 250, 248)
CARD_BG = RGBColor(255, 255, 255)

def set_shape_shadow(shape):
    """Add subtle shadow to shape"""
    spPr = shape._sp.spPr
    shadow = parse_xml(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:outerShdw blurRad="50800" dist="38100" dir="2700000" algn="tl" rotWithShape="0">'
        '<a:srgbClr val="000000"><a:alpha val="23000"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
    )
    spPr.append(shadow)

def add_gradient_fill(shape, color1, color2, angle=90):
    """Add gradient fill to shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2

def create_modern_title_slide(prs):
    """Create a stunning title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full background with gradient effect using shapes
    bg_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5), prs.slide_height)
    bg_left.fill.solid()
    bg_left.fill.fore_color.rgb = PRIMARY
    bg_left.line.fill.background()

    bg_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), 0, Inches(8.33), prs.slide_height)
    bg_right.fill.solid()
    bg_right.fill.fore_color.rgb = WHITE
    bg_right.line.fill.background()

    # Decorative circles
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = PRIMARY_DARK
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(5), Inches(3), Inches(3))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0, 80, 0)
    circle2.line.fill.background()

    # App name on left
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "UBCompass"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Tagline
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(4), Inches(1))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.text = "Navigate Your Campus"
    p.font.size = Pt(24)
    p.font.color.rgb = PRIMARY_LIGHT

    # Right side content
    subtitle = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(7), Inches(0.8))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Campus Navigation App"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.font.bold = True

    main_desc = slide.shapes.add_textbox(Inches(5.5), Inches(2.6), Inches(7), Inches(1.5))
    tf = main_desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "A cross-platform mobile application designed for the University of Buea, providing real-time GPS navigation, indoor wayfinding, and accessibility-aware routing."
    p.font.size = Pt(18)
    p.font.color.rgb = DARK
    p.line_spacing = 1.3

    # Key stats boxes
    stats = [("📱", "Cross-Platform", "iOS & Android"),
             ("📍", "GPS Navigation", "Real-time tracking"),
             ("🏢", "Indoor Maps", "Floor-by-floor")]

    for i, (icon, title_text, desc) in enumerate(stats):
        x = Inches(5.5 + i * 2.5)
        y = Inches(4.5)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.fill.background()

        icon_txt = slide.shapes.add_textbox(x, y + Inches(0.1), Inches(2.2), Inches(0.5))
        tf = icon_txt.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        stat_title = slide.shapes.add_textbox(x, y + Inches(0.55), Inches(2.2), Inches(0.4))
        tf = stat_title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        stat_desc = slide.shapes.add_textbox(x, y + Inches(0.85), Inches(2.2), Inches(0.4))
        tf = stat_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # University badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(6.2), Inches(4), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = PRIMARY
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(Inches(5.5), Inches(6.3), Inches(4), Inches(0.4))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = "University of Buea, Cameroon"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_modern_section(prs, number, title, subtitle=""):
    """Create a modern section divider"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Split background
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(4), prs.slide_height)
    left.fill.solid()
    left.fill.fore_color.rgb = PRIMARY
    left.line.fill.background()

    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), 0, Inches(9.33), prs.slide_height)
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.fill.background()

    # Large number
    num = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(3), Inches(2))
    tf = num.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{number}"
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Section title
    title_box = slide.shapes.add_textbox(Inches(4.5), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(4.5), Inches(3.8), Inches(8), Inches(1))
        tf = sub.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.5), Inches(1.5), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    return slide

def create_problem_solution_slide(prs):
    """Create problem & solution comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Problem & Solution"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Problem card (left)
    prob_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.7), Inches(6), Inches(5))
    prob_card.fill.solid()
    prob_card.fill.fore_color.rgb = RGBColor(255, 235, 235)
    prob_card.line.color.rgb = RGBColor(220, 53, 69)
    prob_card.line.width = Pt(2)

    prob_header = slide.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(5.6), Inches(0.6))
    tf = prob_header.text_frame
    p = tf.paragraphs[0]
    p.text = "❌  The Challenge"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 40, 50)

    problems = [
        "New students get lost finding faculties",
        "No digital navigation exists for UB campus",
        "Indoor building layouts are confusing",
        "No wheelchair-accessible route options",
        "Paper maps are outdated and impractical"
    ]

    prob_list = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(5.4), Inches(4))
    tf = prob_list.text_frame
    tf.word_wrap = True
    for i, prob in enumerate(problems):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {prob}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(14)
        p.line_spacing = 1.4

    # Solution card (right)
    sol_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.7), Inches(6), Inches(5))
    sol_card.fill.solid()
    sol_card.fill.fore_color.rgb = RGBColor(232, 245, 232)
    sol_card.line.color.rgb = PRIMARY
    sol_card.line.width = Pt(2)

    sol_header = slide.shapes.add_textbox(Inches(6.9), Inches(1.9), Inches(5.6), Inches(0.6))
    tf = sol_header.text_frame
    p = tf.paragraphs[0]
    p.text = "✓  UBCompass Solution"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    solutions = [
        "GPS-powered outdoor navigation",
        "SVG indoor floor plans with pathfinding",
        "Real-time 'blue dot' location tracking",
        "Accessibility mode avoids stairs",
        "Accurate ETA with Haversine formula"
    ]

    sol_list = slide.shapes.add_textbox(Inches(7.1), Inches(2.7), Inches(5.4), Inches(4))
    tf = sol_list.text_frame
    tf.word_wrap = True
    for i, sol in enumerate(solutions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {sol}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(14)
        p.line_spacing = 1.4

    return slide

def create_tech_stack_slide(prs):
    """Create modern tech stack visualization"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Tech categories as cards
    categories = [
        ("Frontend", "⚛️", [("React Native", "Cross-platform"), ("Expo", "Dev platform"), ("TypeScript", "Type safety")], Inches(0.4)),
        ("UI/Design", "🎨", [("React Native Paper", "Material Design 3"), ("React Native SVG", "Floor plans")], Inches(3.35)),
        ("Backend", "🗄️", [("Supabase", "PostgreSQL BaaS"), ("OSRM API", "Route calculation")], Inches(6.3)),
        ("Maps", "🗺️", [("OpenStreetMap", "Map tiles"), ("Expo Location", "GPS tracking")], Inches(9.25))
    ]

    for cat_name, icon, techs, x_pos in categories:
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.7), Inches(2.8), Inches(5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(220, 220, 220)

        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(0.9), Inches(1.9), Inches(1), Inches(1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY_LIGHT
        circle.line.fill.background()

        icon_txt = slide.shapes.add_textbox(x_pos + Inches(0.9), Inches(2.1), Inches(1), Inches(0.8))
        tf = icon_txt.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER

        # Category name
        cat_title = slide.shapes.add_textbox(x_pos, Inches(3.05), Inches(2.8), Inches(0.5))
        tf = cat_title.text_frame
        p = tf.paragraphs[0]
        p.text = cat_name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        # Technologies
        y_offset = Inches(3.6)
        for tech_name, tech_desc in techs:
            tech_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos + Inches(0.2), y_offset, Inches(2.4), Inches(0.9))
            tech_box.fill.solid()
            tech_box.fill.fore_color.rgb = LIGHT_BG
            tech_box.line.fill.background()

            tech_title = slide.shapes.add_textbox(x_pos + Inches(0.3), y_offset + Inches(0.1), Inches(2.2), Inches(0.4))
            tf = tech_title.text_frame
            p = tf.paragraphs[0]
            p.text = tech_name
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER

            tech_sub = slide.shapes.add_textbox(x_pos + Inches(0.3), y_offset + Inches(0.45), Inches(2.2), Inches(0.4))
            tf = tech_sub.text_frame
            p = tf.paragraphs[0]
            p.text = tech_desc
            p.font.size = Pt(10)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER

            y_offset += Inches(1)

    return slide

def create_architecture_slide(prs):
    """Create modern architecture diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Mobile App Layer
    app_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.6), Inches(11), Inches(1.2))
    app_box.fill.solid()
    app_box.fill.fore_color.rgb = PRIMARY
    app_box.line.fill.background()

    app_text = slide.shapes.add_textbox(Inches(1), Inches(1.85), Inches(11), Inches(0.7))
    tf = app_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📱  UBCompass Mobile App  |  React Native + Expo + TypeScript"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Screen modules
    screens = [
        ("🗺️ Map", "Outdoor\nNavigation"),
        ("🔍 Search", "Building\nDiscovery"),
        ("🏢 Indoor", "Floor Plan\nNavigation"),
        ("⚙️ Settings", "Preferences &\nAccessibility")
    ]

    for i, (icon_title, desc) in enumerate(screens):
        x = Inches(1.2 + i * 2.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3), Inches(2.5), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = PRIMARY
        box.line.width = Pt(2)

        screen_title = slide.shapes.add_textbox(x, Inches(3.15), Inches(2.5), Inches(0.5))
        tf = screen_title.text_frame
        p = tf.paragraphs[0]
        p.text = icon_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        screen_desc = slide.shapes.add_textbox(x, Inches(3.55), Inches(2.5), Inches(0.7))
        tf = screen_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # State Management
    state_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.5), Inches(11), Inches(0.7))
    state_box.fill.solid()
    state_box.fill.fore_color.rgb = RGBColor(255, 245, 235)
    state_box.line.color.rgb = ACCENT
    state_box.line.width = Pt(2)

    state_text = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11), Inches(0.5))
    tf = state_text.text_frame
    p = tf.paragraphs[0]
    p.text = "📦  Zustand State Management  |  Global State & Preferences"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # External Services
    services = [
        ("📍 Expo Location", "GPS Tracking", ACCENT_BLUE),
        ("🗄️ Supabase", "PostgreSQL Database", RGBColor(62, 207, 142)),
        ("🛣️ OSRM API", "Route Calculation", ACCENT),
        ("🗺️ OpenStreetMap", "Map Tiles", RGBColor(130, 180, 100))
    ]

    for i, (name, desc, color) in enumerate(services):
        x = Inches(1.2 + i * 2.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.5), Inches(2.5), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        svc_title = slide.shapes.add_textbox(x, Inches(5.65), Inches(2.5), Inches(0.5))
        tf = svc_title.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        svc_desc = slide.shapes.add_textbox(x, Inches(6.05), Inches(2.5), Inches(0.4))
        tf = svc_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Connection arrows (simple vertical lines)
    for i in range(4):
        x = Inches(2.45 + i * 2.8)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(5.25), Inches(0.1), Inches(0.25))
        line.fill.solid()
        line.fill.fore_color.rgb = GRAY
        line.line.fill.background()

    return slide

def create_feature_card_slide(prs, title, features):
    """Create a modern feature cards slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Feature cards (2 rows of 2-3)
    for i, (icon, feat_title, description, details) in enumerate(features):
        row = i // 2
        col = i % 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.6 + row * 2.8)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(230, 230, 230)

        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3), Inches(0.9), Inches(0.9))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY_LIGHT
        circle.line.fill.background()

        icon_txt = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.45), Inches(0.9), Inches(0.6))
        tf = icon_txt.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

        # Feature title
        feat_title_box = slide.shapes.add_textbox(x + Inches(1.4), y + Inches(0.35), Inches(4.4), Inches(0.5))
        tf = feat_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = feat_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(1.4), y + Inches(0.8), Inches(4.4), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

        # Details
        details_box = slide.shapes.add_textbox(x + Inches(0.4), y + Inches(1.4), Inches(5.4), Inches(1))
        tf = details_box.text_frame
        tf.word_wrap = True
        for j, detail in enumerate(details):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"✓ {detail}"
            p.font.size = Pt(11)
            p.font.color.rgb = PRIMARY
            p.space_before = Pt(4)

    return slide

def create_algorithm_slide(prs):
    """Create algorithm explanation slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Core Algorithms"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Haversine card
    h_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.6), Inches(6), Inches(2.7))
    h_card.fill.solid()
    h_card.fill.fore_color.rgb = WHITE
    h_card.line.color.rgb = ACCENT_BLUE
    h_card.line.width = Pt(2)

    h_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = h_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📐  Haversine Formula"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    h_desc = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(5.6), Inches(0.5))
    tf = h_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "Calculates distance between GPS coordinates"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY

    h_formula = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(5.6), Inches(1.2))
    tf = h_formula.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "d = 2R × arcsin(√(sin²(Δφ/2) + cos(φ₁)cos(φ₂)sin²(Δλ/2)))"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK
    p = tf.add_paragraph()
    p.text = "Used for: Distance calculation, ETA estimation"
    p.font.size = Pt(11)
    p.font.color.rgb = PRIMARY
    p.space_before = Pt(12)

    # Dijkstra card
    d_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(1.6), Inches(6), Inches(2.7))
    d_card.fill.solid()
    d_card.fill.fore_color.rgb = WHITE
    d_card.line.color.rgb = ACCENT
    d_card.line.width = Pt(2)

    d_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(0.5))
    tf = d_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🛤️  Dijkstra's Algorithm"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    d_desc = slide.shapes.add_textbox(Inches(6.8), Inches(2.3), Inches(5.6), Inches(0.5))
    tf = d_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "Finds shortest path between indoor rooms"
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY

    d_steps = slide.shapes.add_textbox(Inches(6.8), Inches(2.8), Inches(5.6), Inches(1.2))
    tf = d_steps.text_frame
    tf.word_wrap = True
    steps = ["Build graph from room connections", "Process nodes by shortest distance", "Accessibility mode filters stairs"]
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}. {step}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK
        p.space_before = Pt(4)

    # ETA Examples
    eta_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.5), Inches(12.2), Inches(2.1))
    eta_card.fill.solid()
    eta_card.fill.fore_color.rgb = PRIMARY_LIGHT
    eta_card.line.fill.background()

    eta_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.7), Inches(12), Inches(0.5))
    tf = eta_title.text_frame
    p = tf.paragraphs[0]
    p.text = "⏱️  Walking Time Estimation (5 km/h average speed)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    examples = [("100m", "~1 min"), ("500m", "~6 min"), ("1 km", "~12 min"), ("5 km", "~1 hour")]
    for i, (dist, time) in enumerate(examples):
        x = Inches(0.8 + i * 3)

        dist_box = slide.shapes.add_textbox(x, Inches(5.3), Inches(2.5), Inches(0.5))
        tf = dist_box.text_frame
        p = tf.paragraphs[0]
        p.text = dist
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.alignment = PP_ALIGN.CENTER

        time_box = slide.shapes.add_textbox(x, Inches(5.8), Inches(2.5), Inches(0.4))
        tf = time_box.text_frame
        p = tf.paragraphs[0]
        p.text = time
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

        if i < 3:
            arrow = slide.shapes.add_textbox(x + Inches(2.3), Inches(5.5), Inches(0.5), Inches(0.5))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(24)
            p.font.color.rgb = GRAY

    return slide

def create_database_slide(prs):
    """Create database schema slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Database Schema"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Tables
    tables = [
        ("buildings", "🏛️", ["id (UUID)", "name, short_name", "latitude, longitude", "floors, category", "accessible_entrance"], Inches(0.4)),
        ("rooms", "🚪", ["id (UUID)", "building_id (FK)", "name, room_number", "floor, type", "node_x, node_y"], Inches(4.4)),
        ("indoor_graph", "🕸️", ["id (UUID)", "building_id (FK)", "node_x, node_y", "is_stairs, is_ramp", "connections[]"], Inches(8.4))
    ]

    for table_name, icon, fields, x in tables:
        # Table card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(220, 220, 220)

        # Header
        t_header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.8), Inches(0.7))
        t_header.fill.solid()
        t_header.fill.fore_color.rgb = PRIMARY
        t_header.line.fill.background()

        t_title = slide.shapes.add_textbox(x, Inches(1.75), Inches(3.8), Inches(0.5))
        tf = t_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon}  {table_name}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Fields
        fields_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(2.5))
        tf = fields_box.text_frame
        for i, field in enumerate(fields):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {field}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            p.space_before = Pt(8)

    # Campus data summary
    data_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(5.3), Inches(12.2), Inches(1.3))
    data_box.fill.solid()
    data_box.fill.fore_color.rgb = LIGHT_BG
    data_box.line.fill.background()

    data_title = slide.shapes.add_textbox(Inches(0.6), Inches(5.45), Inches(12), Inches(0.4))
    tf = data_title.text_frame
    p = tf.paragraphs[0]
    p.text = "📊  Seeded Campus Data"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    buildings_list = "Main Gate  •  Faculty of Science  •  Faculty of Engineering  •  Administration  •  Library  •  Sports Complex"
    buildings_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.9), Inches(12), Inches(0.5))
    tf = buildings_box.text_frame
    p = tf.paragraphs[0]
    p.text = buildings_list
    p.font.size = Pt(12)
    p.font.color.rgb = DARK

    return slide

def create_demo_overview_slide(prs):
    """Create demo overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Live Demo Walkthrough"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Demo steps
    steps = [
        ("1", "App Launch", "Splash screen → Map view", "30s"),
        ("2", "GPS Permission", "Grant location → Blue dot", "30s"),
        ("3", "Search", "Find Faculty of Science", "1m"),
        ("4", "Directions", "Route + Distance + ETA", "1m"),
        ("5", "Indoor Nav", "Floor plan + Pathfinding", "2m"),
        ("6", "Accessibility", "Toggle → Avoid stairs", "30s")
    ]

    for i, (num, step_title, desc, time) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.7 + row * 2.7)

        # Step card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(220, 220, 220)

        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()

        num_txt = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.3), Inches(0.7), Inches(0.5))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Step title
        step_title_box = slide.shapes.add_textbox(x + Inches(1), y + Inches(0.3), Inches(2.7), Inches(0.5))
        tf = step_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK

        # Description
        desc_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1), Inches(3.5), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY

        # Time badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(2.7), y + Inches(1.9), Inches(1), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()

        time_txt = slide.shapes.add_textbox(x + Inches(2.7), y + Inches(1.97), Inches(1), Inches(0.3))
        tf = time_txt.text_frame
        p = tf.paragraphs[0]
        p.text = time
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def create_roadmap_slide(prs):
    """Create future roadmap slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Development Roadmap"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Timeline line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.5), Inches(12), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()

    phases = [
        ("Completed", "✅", ["Project setup", "Database design", "Outdoor navigation", "Indoor maps", "Search & Settings"], PRIMARY, Inches(0.5)),
        ("Current", "🔄", ["Testing", "Bug fixes", "Performance optimization"], ACCENT, Inches(4.5)),
        ("Next", "🚀", ["More floor plans", "Offline support", "Voice navigation"], ACCENT_BLUE, Inches(8.5))
    ]

    for phase_name, icon, items, color, x in phases:
        # Circle on timeline
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), Inches(3.25), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()

        # Phase card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.5), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)

        phase_title = slide.shapes.add_textbox(x, Inches(1.75), Inches(3.5), Inches(0.5))
        tf = phase_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon}  {phase_name}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        items_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(3.1), Inches(0.8))
        tf = items_box.text_frame
        for i, item in enumerate(items[:3]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = DARK

        # Items below timeline
        below_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4), Inches(3.5), Inches(2.5))
        below_card.fill.solid()
        below_card.fill.fore_color.rgb = LIGHT_BG
        below_card.line.fill.background()

        below_items = slide.shapes.add_textbox(x + Inches(0.2), Inches(4.2), Inches(3.1), Inches(2.2))
        tf = below_items.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"✓ {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = PRIMARY if phase_name == "Completed" else DARK
            p.space_before = Pt(6)

    return slide

def create_thank_you_slide(prs):
    """Create thank you slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full gradient background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Decorative circles
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(5), Inches(5))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = PRIMARY_DARK
    circle1.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0, 80, 0)
    circle2.line.fill.background()

    # Thank you text
    ty_box = slide.shapes.add_textbox(Inches(0), Inches(2), prs.slide_width, Inches(1.5))
    tf = ty_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(3.5), prs.slide_width, Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "UBCompass — Navigating the Future of Campus Mobility"
    p.font.size = Pt(22)
    p.font.color.rgb = PRIMARY_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Credits
    credits_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(4.8), Inches(7), Inches(1.5))
    credits_box.fill.solid()
    credits_box.fill.fore_color.rgb = RGBColor(0, 80, 0)
    credits_box.line.fill.background()

    credits = slide.shapes.add_textbox(Inches(3), Inches(5), Inches(7), Inches(1.2))
    tf = credits.text_frame
    p = tf.paragraphs[0]
    p.text = "Supervisor: Dr. Mougnol Romeo"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Development Team: Group 5"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "University of Buea, Cameroon"
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY_LIGHT
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8)

    return slide

def create_qa_slide(prs):
    """Create Q&A slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Split background
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(6), prs.slide_height)
    left.fill.solid()
    left.fill.fore_color.rgb = PRIMARY
    left.line.fill.background()

    right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6), 0, Inches(7.33), prs.slide_height)
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.fill.background()

    # Q&A on left
    qa_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5), Inches(2))
    tf = qa_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Contact info on right
    contact_title = slide.shapes.add_textbox(Inches(6.5), Inches(2), Inches(6), Inches(0.6))
    tf = contact_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Let's Connect"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(2.7), Inches(1.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    contact_items = [
        ("🎓", "University of Buea"),
        ("👨‍🏫", "Dr. Mougnol Romeo (Supervisor)"),
        ("👥", "Group 5 Development Team"),
        ("📧", "Contact via University Portal")
    ]

    for i, (icon, text) in enumerate(contact_items):
        item_box = slide.shapes.add_textbox(Inches(6.5), Inches(3.2 + i * 0.7), Inches(6), Inches(0.6))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon}  {text}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK

    return slide

def main():
    # Create widescreen presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Title
    create_modern_title_slide(prs)

    # SLIDE 2: Section - Problem & Solution
    create_modern_section(prs, 1, "Problem & Solution", "Understanding the challenge and our approach")

    # SLIDE 3: Problem vs Solution
    create_problem_solution_slide(prs)

    # SLIDE 4: Section - Architecture
    create_modern_section(prs, 2, "System Architecture", "How UBCompass is built")

    # SLIDE 5: Tech Stack
    create_tech_stack_slide(prs)

    # SLIDE 6: Architecture Diagram
    create_architecture_slide(prs)

    # SLIDE 7: Section - Features
    create_modern_section(prs, 3, "Key Features", "What makes UBCompass powerful")

    # SLIDE 8: Navigation Features
    create_feature_card_slide(prs, "Navigation Features", [
        ("📍", "GPS Navigation", "Real-time location tracking with blue dot indicator",
         ["Live GPS with accuracy circle", "OSRM walking routes", "Turn-by-turn directions"]),
        ("🏢", "Indoor Navigation", "SVG floor plans with room-level pathfinding",
         ["Pinch-to-zoom floor plans", "Dijkstra shortest path", "Multi-floor support"]),
        ("♿", "Accessibility Mode", "Wheelchair-friendly routing that avoids stairs",
         ["One-tap toggle", "Ramp-only paths", "Accessible entrances"]),
        ("🔍", "Smart Search", "Instant search across buildings and rooms",
         ["Real-time filtering", "Category icons", "Recent searches"])
    ])

    # SLIDE 9: Section - Algorithms
    create_modern_section(prs, 4, "Core Algorithms", "The math behind navigation")

    # SLIDE 10: Algorithms
    create_algorithm_slide(prs)

    # SLIDE 11: Section - Database
    create_modern_section(prs, 5, "Database Design", "PostgreSQL via Supabase")

    # SLIDE 12: Database Schema
    create_database_slide(prs)

    # SLIDE 13: Section - Demo
    create_modern_section(prs, 6, "Live Demo", "See UBCompass in action")

    # SLIDE 14: Demo Overview
    create_demo_overview_slide(prs)

    # SLIDE 15: Section - Roadmap
    create_modern_section(prs, 7, "Development Roadmap", "Where we are and where we're going")

    # SLIDE 16: Roadmap
    create_roadmap_slide(prs)

    # SLIDE 17: Thank You
    create_thank_you_slide(prs)

    # SLIDE 18: Q&A
    create_qa_slide(prs)

    # Save
    output_path = "/home/daytona/project/UBCompass_Presentation_v2.pptx"
    prs.save(output_path)
    print(f"✅ Professional presentation saved: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return output_path

if __name__ == "__main__":
    main()
